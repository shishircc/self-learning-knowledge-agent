"""Document upload — streaming readers + type classifier.

Supported formats:
  pdf   → pypdf (streaming page-at-a-time)
  docx  → python-docx (paragraphs in order)
  pptx  → python-pptx (one chunk per slide)
  txt / md → buffered read, split into paragraph chunks

The design (see DESIGN.md §"Document ingestion" / TDS §5.11) is that a
document produces a stream of `DocumentChunk`s that each route through the
standard write-path — one PDF paragraph or one slide is the atomic unit
of `new_knowledge`.

For PDFs, the first N pages are sent to the small LM to classify the
document as `word_processing` (prose, split by paragraph) or `presentation`
(one chunk per page). DOCX is always word_processing; PPTX is always
presentation; text/markdown are always word_processing.

Readers are all *lazy*: `read_document` is an async generator that yields
chunks as they become available, so a 100-page PDF doesn't have to be fully
parsed before the agent starts writing packets.
"""
from __future__ import annotations

import json
import os
import re
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import AsyncIterator, Iterator

from . import auth, tracing
from .llm import anthropic_client
from .prompts import extract_json_block


# ---------------------------------------------------------------------------
# Public dataclasses
# ---------------------------------------------------------------------------

@dataclass
class DocumentChunk:
    text: str
    page_number: int
    paragraph_index: int | None  # None for presentation-style (chunk == slide)
    source_filename: str
    document_type: str  # "word_processing" | "presentation"

    def chunk_ref(self) -> str:
        """Stable tag threaded through the LLM prompt so it can name the
        originating chunk in `new_knowledge[i].chunk_ref`.
        """
        if self.paragraph_index is None:
            return f"P{self.page_number}"
        return f"P{self.page_number}.{self.paragraph_index}"

    def to_dict(self) -> dict:
        return asdict(self)


@dataclass
class DocumentMetadata:
    filename: str                     # display name, usually basename
    path: str                          # absolute path on disk
    format: str                        # "pdf" | "docx" | "pptx" | "text" | "markdown"
    document_type: str                 # "word_processing" | "presentation"
    page_count: int | None             # None when unknown at open time
    file_authoritativeness: float
    size_bytes: int


class DocumentReadError(Exception):
    """Raised when a document cannot be opened / decoded / classified."""


class UnsupportedDocumentFormat(DocumentReadError):
    """Raised for formats not in `ingest_doc_allowed_formats`."""


# ---------------------------------------------------------------------------
# Format detection
# ---------------------------------------------------------------------------

_FORMAT_BY_EXT: dict[str, str] = {
    ".pdf":  "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".txt":  "text",
    ".md":   "markdown",
    ".markdown": "markdown",
}


def _sniff_format(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    fmt = _FORMAT_BY_EXT.get(ext)
    if fmt is None:
        raise UnsupportedDocumentFormat(
            f"unsupported extension {ext!r}; supported: {sorted(_FORMAT_BY_EXT)}"
        )
    return fmt


def _implied_document_type(fmt: str) -> str | None:
    """Non-PDF formats have a fixed document_type. Return None for PDF."""
    if fmt in {"docx", "text", "markdown"}:
        return "word_processing"
    if fmt == "pptx":
        return "presentation"
    return None


# ---------------------------------------------------------------------------
# Metadata open — format detection + trust resolution
# ---------------------------------------------------------------------------

def open_metadata(path: str, config: dict) -> DocumentMetadata:
    """Resolve everything we can know without reading the whole file.

    - Sniff format from extension.
    - Reject if not in `config.ingest_doc_allowed_formats`.
    - Reject if file is larger than `config.ingest_doc_max_file_bytes`.
    - Resolve `file_authoritativeness` from config.
    """
    if not os.path.exists(path):
        raise DocumentReadError(f"file not found: {path}")

    fmt = _sniff_format(path)
    allowed = config.get("ingest_doc_allowed_formats") or []
    if allowed and fmt not in allowed:
        raise UnsupportedDocumentFormat(
            f"format {fmt!r} not in ingest_doc_allowed_formats {allowed!r}"
        )

    size = os.path.getsize(path)
    max_size = int(config.get("ingest_doc_max_file_bytes", 25_000_000))
    if size > max_size:
        raise DocumentReadError(
            f"file too large: {size} bytes > {max_size} bytes (ingest_doc_max_file_bytes)"
        )

    file_auth = auth.resolve_file_authoritativeness(path, config)

    # document_type is resolved here for non-PDF; PDF is filled by read_document
    # after the classifier runs.
    return DocumentMetadata(
        filename=os.path.basename(path),
        path=os.path.abspath(path),
        format=fmt,
        document_type=_implied_document_type(fmt) or "",  # "" until classified
        page_count=None,
        file_authoritativeness=file_auth,
        size_bytes=size,
    )


# ---------------------------------------------------------------------------
# Paragraph splitting (word_processing)
# ---------------------------------------------------------------------------

_PARA_SPLIT_RE = re.compile(r"\n\s*\n+")
_SENT_SPLIT_RE = re.compile(r"(?<=[.!?])\s+(?=[A-Z0-9])")


def _split_paragraphs(text: str, min_chars: int, max_chars: int) -> list[str]:
    """Split into paragraphs on blank lines; merge tiny paragraphs forward;
    split oversize paragraphs at sentence boundaries.
    """
    raw_paras = [p.strip() for p in _PARA_SPLIT_RE.split(text or "") if p.strip()]
    if not raw_paras:
        return []

    # Forward-merge tiny paragraphs so single-line bullets don't each become
    # their own chunk.
    merged: list[str] = []
    buf = ""
    for p in raw_paras:
        candidate = f"{buf}\n\n{p}" if buf else p
        if len(candidate) < min_chars:
            buf = candidate
            continue
        merged.append(candidate)
        buf = ""
    if buf:
        if merged:
            merged[-1] = merged[-1] + "\n\n" + buf
        else:
            merged.append(buf)

    # Split oversize paragraphs at sentence boundaries.
    out: list[str] = []
    for p in merged:
        if len(p) <= max_chars:
            out.append(p)
            continue
        sentences = _SENT_SPLIT_RE.split(p)
        chunk = ""
        for s in sentences:
            candidate = f"{chunk} {s}".strip() if chunk else s
            if len(candidate) > max_chars and chunk:
                out.append(chunk)
                chunk = s
            else:
                chunk = candidate
        if chunk:
            out.append(chunk)
    return out


# ---------------------------------------------------------------------------
# Format-specific readers — all yield (page_number, page_text) pairs.
# The word_processing readers get chunked by _split_paragraphs downstream;
# presentation readers yield one slide of text per call, no further split.
# ---------------------------------------------------------------------------

def _read_pdf_pages(path: str, max_pages: int) -> Iterator[tuple[int, str]]:
    try:
        from pypdf import PdfReader  # noqa: F401
    except ImportError as e:  # pragma: no cover — deps installed by pyproject
        raise DocumentReadError(
            "pypdf is required to read PDF files. Install with `pip install pypdf`."
        ) from e

    from pypdf import PdfReader

    reader = PdfReader(path)
    n = len(reader.pages)
    for i in range(min(n, max_pages)):
        page = reader.pages[i]
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        yield (i + 1, text)


def _read_docx_paragraphs(path: str, max_pages: int) -> Iterator[tuple[int, str]]:
    try:
        import docx  # noqa: F401
    except ImportError as e:
        raise DocumentReadError(
            "python-docx is required to read DOCX files. Install with "
            "`pip install python-docx`."
        ) from e

    import docx

    doc = docx.Document(path)
    # DOCX has no true "pages" concept without layout — we treat the whole
    # document as page 1 and let the paragraph splitter produce chunks.
    joined = "\n\n".join(p.text or "" for p in doc.paragraphs)
    yield (1, joined)


def _read_pptx_slides(path: str, max_pages: int) -> Iterator[tuple[int, str]]:
    try:
        import pptx  # noqa: F401
    except ImportError as e:
        raise DocumentReadError(
            "python-pptx is required to read PPTX files. Install with "
            "`pip install python-pptx`."
        ) from e

    from pptx import Presentation

    pres = Presentation(path)
    for i, slide in enumerate(pres.slides):
        if i >= max_pages:
            break
        parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text and text.strip():
                parts.append(text.strip())
        yield (i + 1, "\n\n".join(parts))


def _read_text_pages(path: str) -> Iterator[tuple[int, str]]:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    yield (1, text)


# ---------------------------------------------------------------------------
# PDF document-type classifier
# ---------------------------------------------------------------------------

_CLASSIFIER_SYSTEM = (
    "You are a document classifier. Read a sample of extracted PDF text and "
    "decide whether the document is prose-style (paragraphs, narrative flow, "
    "book / article / report) or slide-style (short bullets, one self-contained "
    "idea per page, deck / presentation). Reply with a single valid JSON object."
)

_CLASSIFIER_PROMPT = """Sample pages from a PDF, in order:

{pages_block}

Task: classify this document as either "word_processing" (prose-style,
book / article / report) or "presentation" (slide deck).

Signals for word_processing: long paragraphs, connected sentences, narrative
prose, references cited inline, headings above dense text.

Signals for presentation: short bullet lists, sparse text per page, big
titles above a handful of lines, phrases rather than sentences.

Output a single valid JSON object only, no fences:
{{
  "document_type": "word_processing" | "presentation",
  "reason": "short reason"
}}"""


async def _classify_pdf(
    sample_pages: list[tuple[int, str]], config: dict
) -> str:
    """Return "word_processing" or "presentation". Falls back to
    word_processing on any error — safer default (over-fragments rather than
    over-collapses).
    """
    if not sample_pages:
        return "word_processing"

    pages_block = "\n\n".join(
        f"--- page {p} ---\n{text[:2000]}"  # cap each sample to keep the call cheap
        for p, text in sample_pages
        if text.strip()
    )
    if not pages_block.strip():
        return "word_processing"

    prompt = _CLASSIFIER_PROMPT.format(pages_block=pages_block)

    with tracing.observation(
        "document_type_classifier",
        as_type="generation",
        model=config["small_lm_model"],
        input=[
            {"role": "system", "content": _CLASSIFIER_SYSTEM},
            {"role": "user", "content": prompt},
        ],
    ) as gen:
        try:
            client = anthropic_client()
            resp = await client.messages.create(
                model=config["small_lm_model"],
                max_tokens=256,
                system=_CLASSIFIER_SYSTEM,
                messages=[{"role": "user", "content": prompt}],
            )
            parts = [
                block.text for block in resp.content
                if getattr(block, "type", None) == "text"
            ]
            raw = "".join(parts)
            tracing.update(gen, output=raw)
        except Exception as e:
            tracing.update(gen, output=f"error: {e}")
            return "word_processing"

    try:
        d = json.loads(extract_json_block(raw))
    except Exception:
        return "word_processing"

    doc_type = (d.get("document_type") or "").strip().lower()
    if doc_type not in {"word_processing", "presentation"}:
        return "word_processing"
    return doc_type


# ---------------------------------------------------------------------------
# read_document — the public streaming generator
# ---------------------------------------------------------------------------

async def read_document(
    metadata: DocumentMetadata, config: dict
) -> AsyncIterator[DocumentChunk]:
    """Yield DocumentChunks lazily as the file is parsed.

    - PDF: classifier runs on the first N pages, then the rest of the pages
      are chunked according to the resolved document_type.
    - DOCX / text / markdown: word_processing paragraphs.
    - PPTX: one chunk per slide (presentation).

    Mutates metadata.document_type + metadata.page_count as it discovers them.
    """
    max_pages = int(config.get("ingest_doc_max_pages", 500))
    min_para = int(config.get("ingest_doc_min_paragraph_chars", 120))
    max_para = int(config.get("ingest_doc_max_paragraph_chars", 1500))
    sample_n = int(config.get("ingest_doc_classifier_sample_pages", 3))
    fmt = metadata.format
    filename = metadata.filename

    with tracing.observation(
        "read_document",
        as_type="span",
        input={"filename": filename, "format": fmt},
        metadata={
            "path": metadata.path,
            "file_authoritativeness": metadata.file_authoritativeness,
        },
    ) as span:
        chunks_emitted = 0

        if fmt == "pdf":
            # Buffer the sample so we can classify before yielding.
            sample: list[tuple[int, str]] = []
            page_stream = _read_pdf_pages(metadata.path, max_pages)
            for page_no, text in page_stream:
                sample.append((page_no, text))
                if len(sample) >= sample_n:
                    break

            doc_type = await _classify_pdf(sample, config)
            metadata.document_type = doc_type

            def _pdf_chained() -> Iterator[tuple[int, str]]:
                for p_and_t in sample:
                    yield p_and_t
                for p_and_t in page_stream:
                    yield p_and_t

            for page_no, text in _pdf_chained():
                if doc_type == "presentation":
                    if text.strip():
                        yield DocumentChunk(
                            text=text.strip(),
                            page_number=page_no,
                            paragraph_index=None,
                            source_filename=filename,
                            document_type=doc_type,
                        )
                        chunks_emitted += 1
                else:
                    paras = _split_paragraphs(text, min_para, max_para)
                    for para_idx, para in enumerate(paras):
                        yield DocumentChunk(
                            text=para,
                            page_number=page_no,
                            paragraph_index=para_idx,
                            source_filename=filename,
                            document_type=doc_type,
                        )
                        chunks_emitted += 1

        elif fmt == "docx":
            metadata.document_type = "word_processing"
            for page_no, text in _read_docx_paragraphs(metadata.path, max_pages):
                paras = _split_paragraphs(text, min_para, max_para)
                for para_idx, para in enumerate(paras):
                    yield DocumentChunk(
                        text=para,
                        page_number=page_no,
                        paragraph_index=para_idx,
                        source_filename=filename,
                        document_type="word_processing",
                    )
                    chunks_emitted += 1

        elif fmt == "pptx":
            metadata.document_type = "presentation"
            for page_no, text in _read_pptx_slides(metadata.path, max_pages):
                if not text.strip():
                    continue
                yield DocumentChunk(
                    text=text.strip(),
                    page_number=page_no,
                    paragraph_index=None,
                    source_filename=filename,
                    document_type="presentation",
                )
                chunks_emitted += 1

        elif fmt in {"text", "markdown"}:
            metadata.document_type = "word_processing"
            for page_no, text in _read_text_pages(metadata.path):
                paras = _split_paragraphs(text, min_para, max_para)
                for para_idx, para in enumerate(paras):
                    yield DocumentChunk(
                        text=para,
                        page_number=page_no,
                        paragraph_index=para_idx,
                        source_filename=filename,
                        document_type="word_processing",
                    )
                    chunks_emitted += 1
        else:
            raise UnsupportedDocumentFormat(f"unhandled format {fmt!r}")

        tracing.update(span, output={
            "chunks_emitted": chunks_emitted,
            "document_type": metadata.document_type,
        })


# ---------------------------------------------------------------------------
# Path detection in a user message
# ---------------------------------------------------------------------------

_QUOTED_PATH_RE = re.compile(r'"([^"]+\.(?:pdf|docx|pptx|txt|md|markdown))"', re.IGNORECASE)
_UNQUOTED_PATH_RE = re.compile(
    r"(?:^|\s)((?:/|~|\./|\.\./)[^\s]+\.(?:pdf|docx|pptx|txt|md|markdown))",
    re.IGNORECASE,
)


def extract_file_paths(text: str) -> list[str]:
    """Regex backstop for `is_upload_request` — pull absolute-ish paths out of
    the user message. Mirrors the URL-regex fallback in extraction.py.

    Recognizes:
      - quoted paths:  "/path/to/file.pdf"
      - unquoted paths that start with `/`, `~`, `./`, `../`
    """
    if not text:
        return []
    out: list[str] = []
    seen: set[str] = set()
    for m in _QUOTED_PATH_RE.finditer(text):
        p = m.group(1).strip()
        if p and p.lower() not in seen:
            seen.add(p.lower())
            out.append(p)
    for m in _UNQUOTED_PATH_RE.finditer(text):
        p = m.group(1).strip().rstrip(".,;:)")
        if p and p.lower() not in seen:
            seen.add(p.lower())
            out.append(p)
    return out


def expand_path(path: str) -> str:
    """Expand `~` and env vars, return an absolute path."""
    return os.path.abspath(os.path.expandvars(os.path.expanduser(path)))
